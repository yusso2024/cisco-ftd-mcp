#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
ACCENT_BLUE = RGBColor(0x00, 0x9E, 0xF7)
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xAA)
ACCENT_PURPLE = RGBColor(0x7C, 0x3A, 0xED)
ACCENT_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
ACCENT_YELLOW = RGBColor(0xFA, 0xCC, 0x15)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xA0, 0xAE, 0xC0)
MID_GRAY = RGBColor(0x33, 0x40, 0x55)
CARD_BG = RGBColor(0x1A, 0x23, 0x3B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(2)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=16, color=WHITE, bold=False, space_before=Pt(6), alignment=PP_ALIGN.LEFT):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = space_before
    p.alignment = alignment
    return p


def add_card(slide, left, top, width, height, fill=CARD_BG, border_color=None):
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill_color=fill, line_color=border_color, line_width=Pt(2))
    card.adjustments[0] = 0.05
    return card


def add_icon_circle(slide, left, top, size, color, label=""):
    circle = add_shape(slide, MSO_SHAPE.OVAL, left, top, size, size, fill_color=color)
    if label:
        circle.text_frame.paragraphs[0].text = label
        circle.text_frame.paragraphs[0].font.size = Pt(int(size / Inches(1) * 14))
        circle.text_frame.paragraphs[0].font.color.rgb = WHITE
        circle.text_frame.paragraphs[0].font.bold = True
        circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        circle.text_frame.word_wrap = False
    return circle


def add_arrow_right(slide, left, top, width, color=ACCENT_BLUE):
    arrow = add_shape(slide, MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.4), fill_color=color)
    return arrow


def add_arrow_down(slide, left, top, height, color=ACCENT_BLUE):
    arrow = add_shape(slide, MSO_SHAPE.DOWN_ARROW, left, top, Inches(0.4), height, fill_color=color)
    return arrow


def gradient_bar(slide, left, top, width, height, color):
    bar = add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, fill_color=color)
    return bar


# =====================================================================
# SLIDE 1 — TITLE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)

gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)
gradient_bar(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), ACCENT_CYAN)

add_icon_circle(slide, Inches(5.6), Inches(1.2), Inches(2), ACCENT_BLUE, "MCP")

add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1.2),
             "Model Context Protocol", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.8),
             "How AI Agents Talk to Your Tools", font_size=24, color=ACCENT_CYAN, bold=False, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.6), Inches(11), Inches(0.6),
             "Case Study: Cisco FTD Security Assessment MCP Server", font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

for i, c in enumerate([ACCENT_BLUE, ACCENT_CYAN, ACCENT_PURPLE]):
    add_shape(slide, MSO_SHAPE.OVAL, Inches(4.5 + i * 1.6), Inches(6.4), Inches(0.15), Inches(0.15), fill_color=c)


# =====================================================================
# SLIDE 2 — WHAT IS MCP?
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "What is MCP?", font_size=36, color=WHITE, bold=True)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05), fill_color=ACCENT_PURPLE)

card = add_card(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.2), border_color=ACCENT_BLUE)
tb = add_text_box(slide, Inches(1.1), Inches(1.8), Inches(5), Inches(2),
                  "The Universal Adapter for AI", font_size=22, color=ACCENT_CYAN, bold=True)
tf = tb.text_frame
add_paragraph(tf, "", font_size=8, color=WHITE)
add_paragraph(tf, "MCP is an open protocol that standardises how AI", font_size=16, color=LIGHT_GRAY)
add_paragraph(tf, "applications connect to external data sources and", font_size=16, color=LIGHT_GRAY)
add_paragraph(tf, "tools  --  like a USB-C port for AI integrations.", font_size=16, color=LIGHT_GRAY)

card2 = add_card(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(2.2), border_color=ACCENT_PURPLE)
tb2 = add_text_box(slide, Inches(7.3), Inches(1.8), Inches(5), Inches(2),
                   "Before MCP", font_size=20, color=ACCENT_ORANGE, bold=True)
tf2 = tb2.text_frame
add_paragraph(tf2, "", font_size=8, color=WHITE)
add_paragraph(tf2, "Every AI + tool combo needed custom glue code.", font_size=15, color=LIGHT_GRAY)
add_paragraph(tf2, "N tools x M AI apps = N x M integrations.", font_size=15, color=LIGHT_GRAY)
add_paragraph(tf2, "", font_size=8, color=WHITE)
add_paragraph(tf2, "After MCP", font_size=20, color=ACCENT_CYAN, bold=True)
add_paragraph(tf2, "Build one MCP server per tool.", font_size=15, color=LIGHT_GRAY)
add_paragraph(tf2, "Every AI client speaks it. N + M. Done.", font_size=15, color=LIGHT_GRAY)

items = [
    ("Tools", "Functions the AI can call (e.g. run_assessment)", ACCENT_BLUE),
    ("Resources", "Data the AI can read (e.g. config files, DB rows)", ACCENT_CYAN),
    ("Prompts", "Pre-built prompt templates for common workflows", ACCENT_PURPLE),
]
for idx, (title, desc, color) in enumerate(items):
    x = Inches(0.8 + idx * 4.1)
    y = Inches(4.3)
    add_card(slide, x, y, Inches(3.8), Inches(2.4), border_color=color)
    add_icon_circle(slide, x + Inches(0.2), y + Inches(0.25), Inches(0.7), color, str(idx + 1))
    add_text_box(slide, x + Inches(1.1), y + Inches(0.3), Inches(2.5), Inches(0.5),
                 title, font_size=20, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.2), Inches(3.3), Inches(1),
                 desc, font_size=14, color=LIGHT_GRAY)


# =====================================================================
# SLIDE 3 — MCP ARCHITECTURE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "MCP Architecture", font_size=36, color=WHITE, bold=True)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2.5), Inches(0.05), fill_color=ACCENT_CYAN)

layers = [
    ("AI Application (Host)", "Claude, OpenCode, Cursor, etc.", ACCENT_PURPLE, Inches(1.8)),
    ("MCP Client", "Built into the host — speaks JSON-RPC", ACCENT_BLUE, Inches(3.2)),
    ("MCP Server", "Your code — exposes tools, resources, prompts", ACCENT_CYAN, Inches(4.6)),
    ("Data Source", "APIs, databases, filesystems, devices", ACCENT_ORANGE, Inches(6.0)),
]

for title, desc, color, y_pos in layers:
    add_card(slide, Inches(2), y_pos, Inches(9.3), Inches(1.0), border_color=color)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(2), y_pos, Inches(0.12), Inches(1.0), fill_color=color)
    add_text_box(slide, Inches(2.4), y_pos + Inches(0.05), Inches(4), Inches(0.5),
                 title, font_size=20, color=color, bold=True)
    add_text_box(slide, Inches(6.5), y_pos + Inches(0.05), Inches(4.5), Inches(0.5),
                 desc, font_size=15, color=LIGHT_GRAY)

for y_start in [Inches(2.8), Inches(4.2), Inches(5.6)]:
    add_shape(slide, MSO_SHAPE.DOWN_ARROW, Inches(6.5), y_start, Inches(0.35), Inches(0.4), fill_color=MID_GRAY)

add_text_box(slide, Inches(0.3), Inches(2.0), Inches(1.5), Inches(4.5),
             "JSON-RPC\nover\nstdio\nor\nSSE", font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.8), Inches(0.04), Inches(5.2), fill_color=MID_GRAY)


# =====================================================================
# SLIDE 4 — HOW MCP COMMUNICATION WORKS
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "How MCP Communication Works", font_size=36, color=WHITE, bold=True)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(3.5), Inches(0.05), fill_color=ACCENT_BLUE)

steps = [
    ("1", "DISCOVER", "Client asks server:\n\"What tools do you have?\"", "Server responds with\ntool list + schemas", ACCENT_PURPLE),
    ("2", "INVOKE", "AI decides to call a tool:\n{\"method\": \"run_assessment\"}", "Server executes the\nfunction locally", ACCENT_BLUE),
    ("3", "RESPOND", "Server returns results\nas structured JSON", "AI interprets and\nexplains to user", ACCENT_CYAN),
]

for idx, (num, title, left_text, right_text, color) in enumerate(steps):
    y = Inches(1.7 + idx * 1.9)

    add_icon_circle(slide, Inches(0.5), y + Inches(0.15), Inches(0.7), color, num)

    add_card(slide, Inches(1.5), y, Inches(4.5), Inches(1.5), border_color=color)
    add_text_box(slide, Inches(1.7), y + Inches(0.1), Inches(4), Inches(0.4),
                 title, font_size=18, color=color, bold=True)
    add_text_box(slide, Inches(1.7), y + Inches(0.5), Inches(4), Inches(0.9),
                 left_text, font_size=14, color=LIGHT_GRAY)

    add_arrow_right(slide, Inches(6.2), y + Inches(0.5), Inches(0.8), color=color)

    add_card(slide, Inches(7.2), y, Inches(5.3), Inches(1.5), border_color=color)
    add_text_box(slide, Inches(7.5), y + Inches(0.3), Inches(4.8), Inches(1),
                 right_text, font_size=14, color=LIGHT_GRAY)

add_text_box(slide, Inches(1.5), Inches(0.9), Inches(4.5), Inches(0.4),
             "CLIENT (AI App)", font_size=13, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(7.2), Inches(0.9), Inches(5.3), Inches(0.4),
             "SERVER (Your Code)", font_size=13, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 5 — OPENCODE + MCP
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_ORANGE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "How OpenCode Uses MCP", font_size=36, color=WHITE, bold=True)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(3), Inches(0.05), fill_color=ACCENT_ORANGE)

add_card(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(3.0), border_color=ACCENT_ORANGE)
add_text_box(slide, Inches(1.0), Inches(1.6), Inches(5.4), Inches(0.4),
             "opencode.json", font_size=20, color=ACCENT_ORANGE, bold=True)

config_text = (
    '{\n'
    '  "mcp": {\n'
    '    "cisco-ftd": {\n'
    '      "type": "local",\n'
    '      "command": [\n'
    '        ".venv/bin/python3",\n'
    '        "server.py"\n'
    '      ]\n'
    '    }\n'
    '  }\n'
    '}'
)
add_text_box(slide, Inches(1.2), Inches(2.1), Inches(5.2), Inches(2.3),
             config_text, font_size=13, color=ACCENT_CYAN)

flow_items = [
    ("OpenCode starts", "Reads opencode.json,\nfinds MCP server entries", ACCENT_PURPLE),
    ("Spawns process", "Runs: python3 server.py\nas a local subprocess", ACCENT_BLUE),
    ("stdio pipe", "Sends JSON-RPC via stdin,\nreads responses from stdout", ACCENT_CYAN),
    ("Tools available", "AI can now call any\n@mcp.tool() function", ACCENT_ORANGE),
]

for idx, (title, desc, color) in enumerate(flow_items):
    y = Inches(1.5 + idx * 1.45)
    x = Inches(7.2)

    add_icon_circle(slide, x, y + Inches(0.15), Inches(0.55), color, str(idx + 1))
    add_text_box(slide, x + Inches(0.7), y + Inches(0.0), Inches(4.5), Inches(0.4),
                 title, font_size=17, color=color, bold=True)
    add_text_box(slide, x + Inches(0.7), y + Inches(0.4), Inches(4.8), Inches(0.8),
                 desc, font_size=13, color=LIGHT_GRAY)

    if idx < len(flow_items) - 1:
        add_shape(slide, MSO_SHAPE.DOWN_ARROW, x + Inches(0.12), y + Inches(0.85), Inches(0.3), Inches(0.45), fill_color=MID_GRAY)

add_card(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.3), border_color=MID_GRAY)
add_text_box(slide, Inches(1.0), Inches(4.9), Inches(4), Inches(0.4),
             "Key Point: 100% Local", font_size=20, color=ACCENT_CYAN, bold=True)

bullets_tb = add_text_box(slide, Inches(1.0), Inches(5.4), Inches(11), Inches(1.5),
                          "", font_size=14, color=LIGHT_GRAY)
tf = bullets_tb.text_frame
tf.paragraphs[0].text = "No HTTP server started  --  MCP server runs as a subprocess, communicates over stdin/stdout"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
add_paragraph(tf, "No ports opened  --  nothing is network-accessible, stays inside the Docker container", font_size=14, color=LIGHT_GRAY, space_before=Pt(10))
add_paragraph(tf, "No cloud dependency  --  your FTD configs and credentials never leave the local machine", font_size=14, color=LIGHT_GRAY, space_before=Pt(10))


# =====================================================================
# SLIDE 6 — CASE STUDY: CISCO FTD
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_RED)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Case Study: Cisco FTD Assessment", font_size=36, color=WHITE, bold=True)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(3.5), Inches(0.05), fill_color=ACCENT_RED)

arch_boxes = [
    ("You (or AI)", Inches(0.4), ACCENT_PURPLE),
    ("OpenCode", Inches(2.3), ACCENT_BLUE),
    ("MCP Server\n(server.py)", Inches(4.2), ACCENT_CYAN),
    ("FMC REST API\nor Config File", Inches(6.1), ACCENT_ORANGE),
    ("FTD\nDevices", Inches(8.0), ACCENT_RED),
]

for label, x, color in arch_boxes:
    add_card(slide, x, Inches(1.6), Inches(1.7), Inches(1.2), border_color=color)
    tb = add_text_box(slide, x + Inches(0.1), Inches(1.7), Inches(1.5), Inches(1.0),
                      label, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)

for x_start in [Inches(2.1), Inches(4.0), Inches(5.9), Inches(7.8)]:
    add_arrow_right(slide, x_start, Inches(2.0), Inches(0.3), color=MID_GRAY)

add_card(slide, Inches(0.4), Inches(3.3), Inches(4.3), Inches(3.8), border_color=ACCENT_CYAN)
add_text_box(slide, Inches(0.6), Inches(3.4), Inches(4), Inches(0.4),
             "FILE MODE", font_size=18, color=ACCENT_CYAN, bold=True)

file_mode_items = [
    "Parse show-run text config",
    "Extract: interfaces, ACLs, NAT,",
    "  VPN, logging, SNMP, NTP, mgmt",
    "Run 11 security check categories",
    "No network needed  --  fully offline",
]
tb = add_text_box(slide, Inches(0.6), Inches(3.9), Inches(3.9), Inches(3),
                  "", font_size=13, color=LIGHT_GRAY)
tf = tb.text_frame
for item in file_mode_items:
    p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
    p.text = item
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(6)

add_card(slide, Inches(5.0), Inches(3.3), Inches(4.3), Inches(3.8), border_color=ACCENT_ORANGE)
add_text_box(slide, Inches(5.2), Inches(3.4), Inches(4), Inches(0.4),
             "LIVE MODE", font_size=18, color=ACCENT_ORANGE, bold=True)

live_mode_items = [
    "Connect to Cisco FMC via REST",
    "Token auth (Basic -> X-auth header)",
    "Pull: devices, policies, rules, IPS",
    "Run assessment against live data",
    "Supports token refresh (30 min TTL)",
]
tb = add_text_box(slide, Inches(5.2), Inches(3.9), Inches(3.9), Inches(3),
                  "", font_size=13, color=LIGHT_GRAY)
tf = tb.text_frame
for item in live_mode_items:
    p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
    p.text = item
    p.font.size = Pt(13)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(6)

add_card(slide, Inches(9.6), Inches(3.3), Inches(3.3), Inches(3.8), border_color=ACCENT_PURPLE)
add_text_box(slide, Inches(9.8), Inches(3.4), Inches(3), Inches(0.4),
             "17 MCP TOOLS", font_size=18, color=ACCENT_PURPLE, bold=True)

tools_items = [
    "connect_fmc",
    "load_config_file",
    "list_devices",
    "list_interfaces",
    "get_access_lists / rules",
    "get_nat_rules",
    "get_vpn_config",
    "get_logging_config",
    "get_snmp_config",
    "get_management_access",
    "run_security_assessment",
]
tb = add_text_box(slide, Inches(9.8), Inches(3.9), Inches(3), Inches(3),
                  "", font_size=12, color=LIGHT_GRAY)
tf = tb.text_frame
for item in tools_items:
    p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
    p.text = item
    p.font.size = Pt(12)
    p.font.color.rgb = LIGHT_GRAY
    p.space_before = Pt(4)


# =====================================================================
# SLIDE 7 — ASSESSMENT FINDINGS
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_RED)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "Security Assessment: What It Catches", font_size=36, color=WHITE, bold=True)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(3.8), Inches(0.05), fill_color=ACCENT_RED)

high_findings = [
    ("Overly Permissive ACL", "permit ip any any"),
    ("Telnet in ACL", "permit tcp any any eq 23"),
    ("Internal Host NAT'd Out", "DB server static NAT to outside"),
    ("Weak VPN Crypto", "DES encryption, MD5 integrity"),
    ("Weak DH Group", "Group 2 (1024-bit) factorable"),
    ("No Syslog Server", "Logs lost on reboot"),
    ("Unrestricted SSH", "0.0.0.0/0 on INSIDE"),
    ("Unrestricted ASDM", "0.0.0.0/0 on INSIDE"),
    ("Default SNMP Strings", "public / private"),
    ("Telnet Enabled", "Cleartext credentials"),
]

medium_findings = [
    ("Unrestricted ICMP", "ICMP any any from outside"),
    ("All-Ports Object Access", "DMZ web->DB on all ports"),
    ("Low Logging Level", "errors only, misses warnings"),
    ("Unauthenticated NTP", "Time-spoofing risk"),
    ("No VPN Idle Timeout", "Abandoned sessions"),
    ("SNMPv2c Only", "No auth/encryption"),
]

add_card(slide, Inches(0.3), Inches(1.5), Inches(3.1), Inches(0.5), border_color=ACCENT_RED)
add_text_box(slide, Inches(0.5), Inches(1.55), Inches(2.8), Inches(0.4),
             f"HIGH  ({len(high_findings)})", font_size=18, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)

for idx, (title, desc) in enumerate(high_findings):
    y = Inches(2.15 + idx * 0.5)
    add_shape(slide, MSO_SHAPE.OVAL, Inches(0.5), y + Inches(0.08), Inches(0.18), Inches(0.18), fill_color=ACCENT_RED)
    add_text_box(slide, Inches(0.85), y, Inches(5.5), Inches(0.45),
                 f"{title}  --  {desc}", font_size=12, color=LIGHT_GRAY)

add_card(slide, Inches(7.0), Inches(1.5), Inches(3.1), Inches(0.5), border_color=ACCENT_YELLOW)
add_text_box(slide, Inches(7.2), Inches(1.55), Inches(2.8), Inches(0.4),
             f"MEDIUM  ({len(medium_findings)})", font_size=18, color=ACCENT_YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

for idx, (title, desc) in enumerate(medium_findings):
    y = Inches(2.15 + idx * 0.5)
    add_shape(slide, MSO_SHAPE.OVAL, Inches(7.2), y + Inches(0.08), Inches(0.18), Inches(0.18), fill_color=ACCENT_YELLOW)
    add_text_box(slide, Inches(7.55), y, Inches(5.5), Inches(0.45),
                 f"{title}  --  {desc}", font_size=12, color=LIGHT_GRAY)

add_card(slide, Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.35))
add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35),
             "All 16 checks derived from CIS Benchmarks, Cisco hardening guides, and NIST SP 800-41 firewall best practices",
             font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 8 — HOW TO USE IT
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "How to Use It", font_size=36, color=WHITE, bold=True)
add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05), fill_color=ACCENT_CYAN)

usage_steps = [
    ("Install", "python3 -m venv .venv && .venv/bin/pip install -r requirements.txt", ACCENT_PURPLE),
    ("Configure", 'Add "cisco-ftd" entry to opencode.json pointing to server.py', ACCENT_BLUE),
    ("Restart", "/restart in OpenCode to load the new MCP server", ACCENT_CYAN),
    ("Use", 'Tell the AI: "Load my FTD config and run a security assessment"', ACCENT_ORANGE),
]

for idx, (title, desc, color) in enumerate(usage_steps):
    y = Inches(1.6 + idx * 1.4)

    add_card(slide, Inches(0.8), y, Inches(11.7), Inches(1.1), border_color=color)
    add_icon_circle(slide, Inches(1.1), y + Inches(0.2), Inches(0.65), color, str(idx + 1))
    add_text_box(slide, Inches(2.0), y + Inches(0.1), Inches(2), Inches(0.4),
                 title, font_size=20, color=color, bold=True)
    add_text_box(slide, Inches(2.0), y + Inches(0.5), Inches(10.2), Inches(0.5),
                 desc, font_size=14, color=LIGHT_GRAY)

add_card(slide, Inches(0.8), Inches(7.25), Inches(11.7), Inches(0.15), fill=ACCENT_CYAN)


# =====================================================================
# SLIDE 9 — SUMMARY
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)
gradient_bar(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), ACCENT_CYAN)

add_icon_circle(slide, Inches(5.7), Inches(1.0), Inches(1.8), ACCENT_BLUE, "MCP")

add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
             "Key Takeaways", font_size=38, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

takeaways = [
    ("MCP = standard protocol", "Build once, works with every AI client", ACCENT_BLUE),
    ("stdio = secure by default", "No ports, no HTTP, no cloud  --  runs locally", ACCENT_CYAN),
    ("Tools = decorated functions", "@mcp.tool() turns any Python function into an AI-callable tool", ACCENT_PURPLE),
    ("Dual mode architecture", "Same server handles live API and offline config parsing", ACCENT_ORANGE),
]

for idx, (title, desc, color) in enumerate(takeaways):
    x = Inches(1.5 + (idx % 2) * 5.5)
    y = Inches(4.0 + (idx // 2) * 1.5)
    add_shape(slide, MSO_SHAPE.OVAL, x, y + Inches(0.08), Inches(0.2), Inches(0.2), fill_color=color)
    add_text_box(slide, x + Inches(0.4), y, Inches(4.5), Inches(0.35),
                 title, font_size=17, color=color, bold=True)
    add_text_box(slide, x + Inches(0.4), y + Inches(0.4), Inches(4.8), Inches(0.5),
                 desc, font_size=14, color=LIGHT_GRAY)

out_path = "/home/opencode/workspace/local-projects/cisco-ftd-mcp/MCP_Cisco_FTD_Assessment.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
